from __future__ import annotations

from pathlib import Path

import pandas as pd
import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
FIGURES = ROOT / "figures"
REPORT = ROOT / "report"
SLIDES = ROOT / "slides"


def fmt(value: object, digits: int = 4) -> str:
    if pd.isna(value):
        return "N/A"
    if isinstance(value, float):
        return f"{value:.{digits}f}"
    return str(value)


def load_results() -> dict[str, pd.DataFrame]:
    return {
        "baseline": pd.read_csv(RESULTS / "baseline_results.csv"),
        "noise": pd.read_csv(RESULTS / "dp_sgd_noise_results.csv"),
        "clip": pd.read_csv(RESULTS / "max_grad_norm_results.csv"),
        "batch": pd.read_csv(RESULTS / "batch_size_results.csv"),
        "multi": pd.read_csv(RESULTS / "multi_seed_results.csv"),
        "mia": pd.read_csv(RESULTS / "mia_results.csv"),
        "summary": pd.read_csv(RESULTS / "final_summary.csv"),
        "tuning": pd.read_csv(RESULTS / "tuning_results.csv")
        if (RESULTS / "tuning_results.csv").exists()
        else pd.DataFrame(),
        "noise_clip": pd.read_csv(RESULTS / "noise_clip_grid_results.csv")
        if (RESULTS / "noise_clip_grid_results.csv").exists()
        else pd.DataFrame(),
        "tuned_baseline": pd.read_csv(RESULTS / "tuned_baseline_results.csv")
        if (RESULTS / "tuned_baseline_results.csv").exists()
        else pd.DataFrame(),
        "strong_baseline": pd.read_csv(RESULTS / "strong_baseline_full_train.csv")
        if (RESULTS / "strong_baseline_full_train.csv").exists()
        else pd.DataFrame(),
        "boosted_baseline": pd.read_csv(RESULTS / "boosted_feature_baseline_results.csv")
        if (RESULTS / "boosted_feature_baseline_results.csv").exists()
        else pd.DataFrame(),
        "model_grid": pd.read_csv(RESULTS / "model_clip_noise_grid_results.csv")
        if (RESULTS / "model_clip_noise_grid_results.csv").exists()
        else pd.DataFrame(),
    }


def markdown_table(df: pd.DataFrame, columns: list[str]) -> str:
    header = "| " + " | ".join(columns) + " |"
    sep = "| " + " | ".join(["---"] * len(columns)) + " |"
    rows = []
    for _, row in df.iterrows():
        rows.append("| " + " | ".join(fmt(row[col]) for col in columns) + " |")
    return "\n".join([header, sep, *rows])


def create_report_md(results: dict[str, pd.DataFrame]) -> str:
    baseline = results["baseline"]
    noise = results["noise"]
    clip = results["clip"]
    batch = results["batch"]
    multi = results["multi"]
    mia = results["mia"]
    tuning = results["tuning"]
    noise_clip = results["noise_clip"]
    tuned_baseline = results["tuned_baseline"]
    strong_baseline = results["strong_baseline"]
    boosted_baseline = results["boosted_baseline"]
    model_grid = results["model_grid"]

    balanced = noise[noise["noise_multiplier"] == 1.5].iloc[0]
    strongest = noise.sort_values("epsilon").iloc[0]
    best_clip = clip.sort_values("f1_score", ascending=False).iloc[0]
    best_tuned = None if tuning.empty else tuning[tuning["status"] == "ok"].sort_values("f1_score", ascending=False).iloc[0]
    tuning_section = ""
    if best_tuned is not None:
        tuning_section = f"""
## 9. GPU Tuning Result

The tuning run used the CUDA PyTorch environment on the NVIDIA T1200 Laptop GPU and focused on improving DP-SGD with `noise_multiplier = 1.5`. The best configuration found was:

| max_grad_norm | batch_size | epochs | learning_rate | schedule | epsilon | accuracy | F1-score |
|---:|---:|---:|---:|---|---:|---:|---:|
| {fmt(best_tuned["max_grad_norm"])} | {fmt(best_tuned["batch_size"], 0)} | {fmt(best_tuned["epochs"], 0)} | {fmt(best_tuned["learning_rate"])} | {best_tuned["schedule"]} | {fmt(best_tuned["epsilon"])} | {fmt(best_tuned["accuracy"])} | {fmt(best_tuned["f1_score"])} |

The search stopped because later attempts changed F1 by less than the predefined meaningful threshold or required a larger epsilon without a clear utility gain. Detailed configuration notes are stored in `docs/tuning_notes.md`.
"""
    strong_baseline_section = ""
    if not strong_baseline.empty:
        best_strong = strong_baseline.iloc[0]
        strong_baseline_section = f"""
## 5.1 Strong Non-DP Baseline

After tuning the simple MLP baseline, a stronger classical baseline was tested to check whether the project had an under-powered non-DP reference. The strongest clean local result came from HistGradientBoosting trained on the full Adult training file, using a threshold selected from an earlier training-validation split.

| model | default accuracy | default F1 | tuned accuracy | tuned F1 | ROC-AUC | PR-AUC |
|---|---:|---:|---:|---:|---:|---:|
| {best_strong["model_name"]} | {fmt(best_strong["default_accuracy"])} | {fmt(best_strong["default_f1_score"])} | {fmt(best_strong["accuracy"])} | {fmt(best_strong["f1_score"])} | {fmt(best_strong["roc_auc"])} | {fmt(best_strong["pr_auc"])} |

This is now the strongest clean non-DP reference. It improves the baseline, but it still does not reach 0.90 accuracy or 0.90 F1 on the original Adult test split.
"""
    boosted_section = ""
    if not boosted_baseline.empty:
        best_boosted = boosted_baseline[boosted_baseline["status"] == "ok"].sort_values("f1_score", ascending=False).iloc[0]
        boosted_section = f"""
## 5.2 Boosted Feature Baseline

XGBoost, LightGBM, CatBoost, and HistGradientBoosting were tested with additional binning and interaction features. Thresholds were selected on a training-validation split, then the selected model was evaluated on the held-out Adult test file.

| model | default accuracy | default F1 | tuned accuracy | tuned F1 | ROC-AUC | PR-AUC |
|---|---:|---:|---:|---:|---:|---:|
| {best_boosted["model_name"]} | {fmt(best_boosted["default_accuracy"])} | {fmt(best_boosted["default_f1_score"])} | {fmt(best_boosted["accuracy"])} | {fmt(best_boosted["f1_score"])} | {fmt(best_boosted["roc_auc"])} | {fmt(best_boosted["pr_auc"])} |

The best boosted/feature-engineered baseline still does not reach 0.90 accuracy or 0.90 F1, so a 0.90 target should be treated as a stretch goal or leak-check trigger.
"""
    model_grid_section = ""
    if not model_grid.empty:
        ok_grid = model_grid[model_grid["status"] == "ok"]
        best_grid = ok_grid.sort_values("f1_score", ascending=False).iloc[0]
        worst_grid = ok_grid.sort_values("f1_score", ascending=True).iloc[0]
        best_by_noise = (
            ok_grid.sort_values("f1_score", ascending=False)
            .groupby("noise_multiplier", as_index=False)
            .first()
            .sort_values("noise_multiplier")
        )
        model_grid_section = f"""
## 10.1 Model x Clipping x Noise Grid

The full DP grid evaluates four MLP variants across the same clipping and noise ranges, for 144 GPU runs.

| noise_multiplier | model_variant | max_grad_norm | epsilon | F1-score |
|---:|---|---:|---:|---:|
""" + "\n".join(
            f"| {fmt(row['noise_multiplier'], 1)} | {row['model_variant']} | {fmt(row['max_grad_norm'], 1)} | {fmt(row['epsilon'])} | {fmt(row['f1_score'])} |"
            for _, row in best_by_noise.iterrows()
        ) + f"""

Best DP grid result: `{best_grid["model_variant"]}`, noise `{fmt(best_grid["noise_multiplier"], 1)}`, clip `{fmt(best_grid["max_grad_norm"], 1)}`, F1 `{fmt(best_grid["f1_score"])}`. Worst grid result: `{worst_grid["model_variant"]}`, noise `{fmt(worst_grid["noise_multiplier"], 1)}`, clip `{fmt(worst_grid["max_grad_norm"], 1)}`, F1 `{fmt(worst_grid["f1_score"])}`.
"""
    stress_section = ""
    if not noise_clip.empty:
        baseline_f1 = None
        if not tuned_baseline.empty:
            baseline_f1 = float(tuned_baseline.sort_values("f1_score", ascending=False).iloc[0]["f1_score"])
        best_by_noise = (
            noise_clip[noise_clip["status"] == "ok"]
            .sort_values("f1_score", ascending=False)
            .groupby("noise_multiplier", as_index=False)
            .first()
            .sort_values("noise_multiplier")
        )
        first_drop = "N/A"
        if baseline_f1 is not None:
            drop_rows = best_by_noise[best_by_noise["f1_score"] <= baseline_f1 - 0.05]
            if not drop_rows.empty:
                first_drop = fmt(drop_rows.iloc[0]["noise_multiplier"], 1)
        stress_section = f"""
## 10. Noise x Clipping Stress Test

This additional GPU experiment was added after feedback that the original noise range may not be large enough to show clear utility collapse. It uses F1-score as the primary metric and evaluates a joint grid of `noise_multiplier x max_grad_norm`. The comparison reference is the tuned non-DP MLP baseline when available, so low-noise DP is not allowed to look better only because the original baseline was under-tuned.

{markdown_table(best_by_noise, ["noise_multiplier", "max_grad_norm", "epsilon", "precision", "recall", "f1_score"])}

Tuned non-DP baseline F1: {fmt(baseline_f1) if baseline_f1 is not None else "N/A"}. The best DP result in the grid is below this tuned baseline. Best-per-noise F1 first drops by at least 0.05 below tuned baseline at noise `{first_drop}`, and larger stress-test noise values such as 10.0, 15.0, and 20.0 show a much clearer utility drop. Detailed notes are stored in `docs/noise_clip_stress_notes.md`.
"""

    return f"""# Privacy-Preserving Data with Differential Privacy using DP-SGD

## 1. Introduction

This project studies whether DP-SGD can be integrated into a normal PyTorch training pipeline to reduce privacy risk while preserving useful performance on the UCI Adult personal tabular dataset.

## 2. Threat Model

The assumed attacker can query or inspect a trained model and may try to infer whether a specific person's record was part of the training data. DP-SGD limits the influence of any individual record through per-example gradient clipping, Gaussian noise, and privacy accounting. The project also evaluates a simple confidence-based membership inference attack, but it does not claim protection against every possible privacy attack.

## 3. Dataset and Method

- Dataset: UCI Adult / Census Income.
- Task: binary classification, `>50K` vs `<=50K`.
- Split: original `adult.data` train file and `adult.test` test file.
- Preprocessing: drop rows with `?`, standardize numeric columns, one-hot encode categorical columns.
- Models: Logistic Regression baseline, non-DP MLP, and DP-SGD MLP with Opacus.
- Main DP parameters: `delta = 1e-5`, `epochs = 20`, `learning_rate = 0.05`, `momentum = 0.9`.

## 4. Baseline Results

{markdown_table(baseline, ["model_name", "accuracy", "precision", "recall", "f1_score", "roc_auc", "pr_auc", "training_time"])}

## 5. Noise Multiplier Sweep

{markdown_table(noise, ["noise_multiplier", "epsilon", "accuracy", "precision", "recall", "f1_score", "roc_auc", "pr_auc", "training_time"])}

The balanced configuration remains `noise_multiplier = 1.5`: epsilon = {fmt(balanced["epsilon"])}, accuracy = {fmt(balanced["accuracy"])}, and F1-score = {fmt(balanced["f1_score"])}. The strongest formal privacy setting in this sweep is `noise_multiplier = {fmt(strongest["noise_multiplier"], 1)}` with epsilon = {fmt(strongest["epsilon"])}.

{strong_baseline_section}
{boosted_section}

## 6. Max Grad Norm Sweep

{markdown_table(clip, ["max_grad_norm", "epsilon", "accuracy", "precision", "recall", "f1_score", "roc_auc", "pr_auc", "training_time"])}

With fixed noise multiplier, sample rate, epochs, and delta, epsilon is unchanged across clipping norms. The best F1-score in this run is obtained at `max_grad_norm = {fmt(best_clip["max_grad_norm"], 1)}`.

## 7. Batch Size Sweep

{markdown_table(batch, ["batch_size", "epsilon", "accuracy", "precision", "recall", "f1_score", "roc_auc", "pr_auc", "training_time"])}

Batch size changes privacy accounting through the sampling rate. In this run, batch size 64 gives the smallest epsilon but has much weaker F1-score, while batch size 512 has stronger utility but a larger epsilon.

## 8. Multi-Seed Robustness

{markdown_table(multi, ["model_config", "mean_accuracy", "std_accuracy", "mean_f1_score", "std_f1_score", "mean_roc_auc", "std_roc_auc", "mean_epsilon", "std_epsilon"])}

The standard deviations are small, which suggests that the main results are reasonably stable across the selected seeds.

## 9. Membership Inference Attack

{markdown_table(mia, ["model_name", "epsilon", "attack_auc", "mean_train_confidence", "mean_test_confidence", "confidence_gap"])}

All attack AUC values are close to 0.5. This means the simple confidence-based attack is close to random guessing in this setup. The result should be interpreted carefully: it is an empirical attack check, while epsilon remains the formal privacy metric.

{tuning_section}
{stress_section}
{model_grid_section}

## 11. Comparison with Abadi et al.

| Criterion | Abadi et al. | This project |
|---|---|---|
| Dataset | MNIST, CIFAR-10 | UCI Adult |
| Data type | Images | Personal tabular data |
| Model | Neural network / CNN | Logistic Regression and MLP |
| Privacy mechanism | DP-SGD | DP-SGD through Opacus |
| Privacy accounting | Moments Accountant | Opacus PrivacyEngine |
| Analysis axis | Epsilon, delta, accuracy | Epsilon, utility metrics, training time, MIA AUC |

## 12. Limitations

- The membership inference attack is simple and confidence-based.
- The project uses one tabular dataset and a small MLP.
- The final stress-test and tuning runs use the local NVIDIA T1200 GPU, but the project is still limited by laptop hardware.
- Secure RNG was off in Opacus for faster experimentation; production-grade privacy training should enable secure mode.

## 13. Conclusion

The final pipeline shows that DP-SGD can be integrated into model training for a personal tabular dataset. Increasing noise reduces epsilon, and the new stress test shows that F1-score drops clearly once the noise multiplier is pushed into much larger values. The balanced setting `noise_multiplier = 1.5` remains a useful privacy-utility compromise, while extreme noise settings are useful for demonstrating utility collapse. The membership inference attack results show low confidence-based membership signal in this particular experiment.
"""


def add_docx_table(document: Document, df: pd.DataFrame, columns: list[str]) -> None:
    table = document.add_table(rows=1, cols=len(columns))
    table.style = "Table Grid"
    for i, col in enumerate(columns):
        table.rows[0].cells[i].text = col
    for _, row in df.iterrows():
        cells = table.add_row().cells
        for i, col in enumerate(columns):
            cells[i].text = fmt(row[col])


def create_docx(results: dict[str, pd.DataFrame]) -> None:
    document = Document()
    document.add_heading("Privacy-Preserving Data with Differential Privacy using DP-SGD", 0)
    document.add_paragraph("Final report generated from saved CSV results.")

    sections = [
        ("Baseline Results", results["baseline"], ["model_name", "accuracy", "f1_score", "roc_auc", "pr_auc", "training_time"]),
        ("Noise Multiplier Sweep", results["noise"], ["noise_multiplier", "epsilon", "accuracy", "f1_score", "roc_auc", "pr_auc"]),
        ("Max Grad Norm Sweep", results["clip"], ["max_grad_norm", "epsilon", "accuracy", "f1_score", "roc_auc", "pr_auc"]),
        ("Batch Size Sweep", results["batch"], ["batch_size", "epsilon", "accuracy", "f1_score", "roc_auc", "pr_auc"]),
        ("Multi-Seed Robustness", results["multi"], ["model_config", "mean_accuracy", "std_accuracy", "mean_f1_score", "std_f1_score", "mean_epsilon"]),
        ("Membership Inference Attack", results["mia"], ["model_name", "epsilon", "attack_auc", "confidence_gap"]),
    ]

    for title, df, columns in sections:
        document.add_heading(title, 1)
        add_docx_table(document, df, columns)

    tuning = results["tuning"]
    noise_clip = results["noise_clip"]
    tuned_baseline = results["tuned_baseline"]
    strong_baseline = results["strong_baseline"]
    boosted_baseline = results["boosted_baseline"]
    model_grid = results["model_grid"]
    if not tuning.empty:
        tuning_ok = tuning[tuning["status"] == "ok"].copy()
        best_tuned = tuning_ok.sort_values("f1_score", ascending=False).iloc[0]
        document.add_heading("GPU Tuning Result", 1)
        document.add_paragraph(
            "The best tuned DP-SGD configuration found on GPU used "
            f"max_grad_norm={fmt(best_tuned['max_grad_norm'])}, "
            f"batch_size={fmt(best_tuned['batch_size'], 0)}, "
            f"epochs={fmt(best_tuned['epochs'], 0)}, "
            f"learning_rate={fmt(best_tuned['learning_rate'])}, "
            f"schedule={best_tuned['schedule']}. "
            f"It reached epsilon={fmt(best_tuned['epsilon'])}, "
            f"accuracy={fmt(best_tuned['accuracy'])}, and "
            f"F1-score={fmt(best_tuned['f1_score'])}."
        )
    if not strong_baseline.empty:
        document.add_heading("Strong Non-DP Baseline", 1)
        add_docx_table(
            document,
            strong_baseline,
            ["model_name", "default_accuracy", "default_f1_score", "accuracy", "f1_score", "roc_auc", "pr_auc"],
        )
        document.add_paragraph(
            "This stronger clean baseline improves F1, but still does not reach 0.90 accuracy or 0.90 F1 on the original Adult test split."
        )
    if not boosted_baseline.empty:
        best_boosted = boosted_baseline[boosted_baseline["status"] == "ok"].sort_values("f1_score", ascending=False).head(5)
        document.add_heading("Boosted Feature Baseline", 1)
        add_docx_table(
            document,
            best_boosted,
            ["model_name", "default_accuracy", "default_f1_score", "accuracy", "f1_score", "roc_auc", "pr_auc"],
        )
    if not model_grid.empty:
        ok_grid = model_grid[model_grid["status"] == "ok"]
        best_by_model = ok_grid.sort_values("f1_score", ascending=False).groupby("model_variant", as_index=False).first()
        document.add_heading("Model x Clipping x Noise Grid", 1)
        add_docx_table(
            document,
            best_by_model,
            ["model_variant", "noise_multiplier", "max_grad_norm", "epsilon", "accuracy", "f1_score"],
        )
    if not noise_clip.empty:
        baseline_f1 = None
        if not tuned_baseline.empty:
            baseline_f1 = float(tuned_baseline.sort_values("f1_score", ascending=False).iloc[0]["f1_score"])
        best_by_noise = (
            noise_clip[noise_clip["status"] == "ok"]
            .sort_values("f1_score", ascending=False)
            .groupby("noise_multiplier", as_index=False)
            .first()
            .sort_values("noise_multiplier")
        )
        document.add_heading("Noise x Clipping Stress Test", 1)
        add_docx_table(
            document,
            best_by_noise,
            ["noise_multiplier", "max_grad_norm", "epsilon", "precision", "recall", "f1_score"],
        )
        document.add_paragraph(
            "The stress test uses the tuned non-DP MLP as the comparison reference. "
            f"Tuned baseline F1 is {fmt(baseline_f1) if baseline_f1 is not None else 'N/A'}, "
            "and the best DP result in the grid is lower than that reference. Larger noise values make "
            "the F1-score drop much more clearly."
        )

    document.add_heading("Key Interpretation", 1)
    document.add_paragraph(
        "Noise multiplier 1.5 remains a balanced configuration for discussion. "
        "Noise multiplier 3.0 has the smallest epsilon in the sweep. "
        "The confidence-based membership inference attack has AUC values close to 0.5, "
        "so the simple attack is near random guessing in this setup."
    )
    document.save(REPORT / "report.docx")


def add_slide(prs: Presentation, title: str, bullets: list[str], image: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
    if image and image.exists():
        slide.shapes.add_picture(str(image), PptInches(6.2), PptInches(1.55), width=PptInches(3.4))


def create_slides(results: dict[str, pd.DataFrame]) -> None:
    noise = results["noise"]
    balanced = noise[noise["noise_multiplier"] == 1.5].iloc[0]
    strongest = noise.sort_values("epsilon").iloc[0]
    mia = results["mia"]
    tuning = results["tuning"]
    noise_clip = results["noise_clip"]
    tuned_baseline = results["tuned_baseline"]
    strong_baseline = results["strong_baseline"]
    boosted_baseline = results["boosted_baseline"]
    model_grid = results["model_grid"]
    tuned = None if tuning.empty else tuning[tuning["status"] == "ok"].sort_values("f1_score", ascending=False).iloc[0]
    tuned_baseline_best = (
        None if tuned_baseline.empty else tuned_baseline.sort_values("f1_score", ascending=False).iloc[0]
    )
    strong_baseline_best = None if strong_baseline.empty else strong_baseline.iloc[0]
    boosted_best = (
        None if boosted_baseline.empty else boosted_baseline[boosted_baseline["status"] == "ok"].sort_values("f1_score", ascending=False).iloc[0]
    )
    model_grid_best = (
        None if model_grid.empty else model_grid[model_grid["status"] == "ok"].sort_values("f1_score", ascending=False).iloc[0]
    )
    model_grid_worst = (
        None if model_grid.empty else model_grid[model_grid["status"] == "ok"].sort_values("f1_score", ascending=True).iloc[0]
    )
    stress_best_noise = None
    if not noise_clip.empty:
        stress_best_noise = (
            noise_clip[noise_clip["status"] == "ok"]
            .sort_values("f1_score", ascending=False)
            .groupby("noise_multiplier", as_index=False)
            .first()
            .sort_values("noise_multiplier")
        )

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Privacy-Preserving Data with Differential Privacy"
    title_slide.placeholders[1].text = "DP-SGD on UCI Adult Dataset"

    slide_specs = [
        ("Motivation", ["Machine learning models may memorize training records", "Personal data needs privacy in the training pipeline"]),
        ("Problem", ["Attacker may infer whether a person was in training data", "UCI Adult contains personal tabular attributes"]),
        ("Differential Privacy", ["Smaller epsilon means stronger privacy", "Delta is fixed at 1e-5 in this project"]),
        ("DP-SGD Workflow", ["Per-example gradients", "L2 clipping", "Gaussian noise", "Privacy accounting"]),
        ("Source Paper", ["Abadi et al.: Deep Learning with Differential Privacy", "Original benchmarks: MNIST and CIFAR-10"]),
        ("Dataset", ["UCI Adult / Census Income", "30,162 train rows and 15,060 test rows after cleaning", "104 processed input features"]),
        ("Threat Model", ["Attacker can query trained model", "Attack goal: infer training membership", "Defense: bound influence of one record"]),
        ("Experimental Pipeline", ["Preprocess data", "Train baseline models", "Run DP-SGD sweeps", "Evaluate utility and MIA risk"]),
        ("Baseline Models", ["Logistic Regression accuracy = 0.8480", "MLP baseline accuracy = 0.8487", "MLP baseline F1 = 0.6645"]),
        (
            "Strong Baseline Check",
            [
                "Tuned classical baseline: HistGradientBoosting",
                f"Default accuracy = {fmt(strong_baseline_best['default_accuracy']) if strong_baseline_best is not None else 'N/A'}",
                f"Threshold-tuned F1 = {fmt(strong_baseline_best['f1_score']) if strong_baseline_best is not None else 'N/A'}",
                "Clean 0.90 target was not reached",
            ],
        ),
        (
            "Boosted Feature Baseline",
            [
                "Tried XGBoost, LightGBM, CatBoost, HistGradientBoosting",
                "Added numeric bins and interaction features",
                f"Best boosted F1 = {fmt(boosted_best['f1_score']) if boosted_best is not None else 'N/A'}",
                f"Best boosted default accuracy = {fmt(boosted_best['default_accuracy']) if boosted_best is not None else 'N/A'}",
            ],
        ),
        ("Noise Sweep", [f"Balanced noise 1.5 epsilon = {fmt(balanced['epsilon'])}", f"Accuracy = {fmt(balanced['accuracy'])}", f"F1 = {fmt(balanced['f1_score'])}"], FIGURES / "noise_vs_epsilon.png"),
        ("Privacy Utility Tradeoff", ["Higher noise reduces epsilon", "Utility changes are modest on UCI Adult"], FIGURES / "privacy_utility_tradeoff.png"),
        ("Max Grad Norm", ["Epsilon unchanged when noise/sample rate/epochs are fixed", "Best F1 in final run: max_grad_norm = 2.0"], FIGURES / "max_grad_norm_vs_f1.png"),
        ("Batch Size", ["Batch size changes sampling rate", "Batch 64 has lowest epsilon but much weaker F1", "Batch 512 has stronger utility but higher epsilon"], FIGURES / "batch_size_vs_epsilon.png"),
        ("Multi-Seed Robustness", ["Small standard deviations across seeds", "Noise 1.5 remains stable enough for discussion"], FIGURES / "multi_seed_errorbar_accuracy.png"),
        ("Membership Inference Setup", ["Confidence-based attack", "Train samples label 1, test samples label 0", "Metric: attack AUC"]),
        ("Membership Inference Result", [f"Baseline attack AUC = {fmt(mia.iloc[0]['attack_auc'])}", "DP-SGD attack AUC values are also near 0.5", "Simple attack is close to random guessing"], FIGURES / "attack_auc_comparison.png"),
        ("Recommended Configuration", [f"Balanced: noise_multiplier = 1.5", f"Epsilon = {fmt(balanced['epsilon'])}", f"Strongest privacy: noise = {fmt(strongest['noise_multiplier'], 1)}, epsilon = {fmt(strongest['epsilon'])}"]),
        (
            "Tuned DP-SGD",
            [
                "GPU tuning focused on noise multiplier 1.5",
                f"Best F1 = {fmt(tuned['f1_score']) if tuned is not None else 'N/A'}",
                f"clip={fmt(tuned['max_grad_norm']) if tuned is not None else 'N/A'}, batch={fmt(tuned['batch_size'], 0) if tuned is not None else 'N/A'}",
                f"epochs={fmt(tuned['epochs'], 0) if tuned is not None else 'N/A'}, lr={fmt(tuned['learning_rate']) if tuned is not None else 'N/A'}",
            ],
            FIGURES / "tuning_f1_by_config.png",
        ),
        (
            "Noise Stress Test",
            [
                "Added joint grid: noise x clipping norm",
                "Metric focus: F1-score",
                f"Tuned non-DP F1 = {fmt(tuned_baseline_best['f1_score']) if tuned_baseline_best is not None else 'N/A'}",
                "Best DP grid result is below tuned baseline",
                "Noise 10-20 shows clear utility drop",
                f"Best F1 at noise 20 = {fmt(stress_best_noise[stress_best_noise['noise_multiplier'] == 20.0].iloc[0]['f1_score']) if stress_best_noise is not None and not stress_best_noise[stress_best_noise['noise_multiplier'] == 20.0].empty else 'N/A'}",
            ],
            FIGURES / "extreme_noise_best_f1.png",
        ),
        (
            "Model x Clip x Noise",
            [
                "Completed 144 DP-SGD GPU runs",
                f"Best DP grid F1 = {fmt(model_grid_best['f1_score']) if model_grid_best is not None else 'N/A'}",
                f"Best DP model = {model_grid_best['model_variant'] if model_grid_best is not None else 'N/A'}",
                f"Worst DP grid F1 = {fmt(model_grid_worst['f1_score']) if model_grid_worst is not None else 'N/A'}",
            ],
            FIGURES / "model_clip_noise_best_f1.png",
        ),
        ("Conclusion", ["DP-SGD can be integrated into PyTorch training", "Privacy and utility should be reported together", "MIA AUC complements but does not replace epsilon"]),
    ]

    for spec in slide_specs:
        title, bullets, *image = spec
        add_slide(prs, title, bullets, image[0] if image else None)

    prs.save(SLIDES / "dp_sgd_final_presentation.pptx")


def main() -> None:
    REPORT.mkdir(exist_ok=True)
    SLIDES.mkdir(exist_ok=True)
    results = load_results()
    tuning = results["tuning"]
    if not tuning.empty:
        tuning_ok = tuning[tuning["status"] == "ok"].copy().reset_index(drop=True)
        tuning_ok["config"] = [
            f"{int(row.batch_size)}bs/{row.max_grad_norm}clip/{int(row.epochs)}ep/{row.learning_rate}lr/{row.schedule}"
            for row in tuning_ok.itertuples()
        ]
        plt.figure(figsize=(11, 5))
        plt.plot(range(len(tuning_ok)), tuning_ok["f1_score"], marker="o")
        plt.xticks(range(len(tuning_ok)), tuning_ok["config"], rotation=45, ha="right", fontsize=7)
        plt.ylabel("F1-score")
        plt.title("DP-SGD GPU Tuning F1 by Configuration")
        plt.grid(True)
        plt.tight_layout()
        plt.savefig(FIGURES / "tuning_f1_by_config.png", dpi=200)
        plt.close()
    report_md = create_report_md(results)
    (REPORT / "report.md").write_text(report_md, encoding="utf-8")
    create_docx(results)
    create_slides(results)
    print("Saved:", REPORT / "report.md")
    print("Saved:", REPORT / "report.docx")
    print("Saved:", SLIDES / "dp_sgd_final_presentation.pptx")


if __name__ == "__main__":
    main()
